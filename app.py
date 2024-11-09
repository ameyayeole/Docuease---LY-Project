import streamlit as st
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_google_genai import GoogleGenerativeAIEmbeddings
import google.generativeai as genai
from langchain_community.vectorstores import FAISS
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.chains.question_answering import load_qa_chain
from langchain.prompts import PromptTemplate
from langchain.chains import LLMChain
from dotenv import load_dotenv
from gtts import gTTS
from googletrans import Translator
import io
import os

load_dotenv()
api_key = os.getenv("GOOGLE_API_KEY")
genai.configure(api_key=api_key)

translator = Translator()

LANGUAGES = {
    "en": "English",
    "es": "Spanish",
    "fr": "French",
    "de": "German",
    "zh-cn": "Chinese (Simplified)",
    "hi": "Hindi"
}

def extract_text_from_pdf(pdf_docs):
    text = ""
    for pdf in pdf_docs:
        pdf_reader = PdfReader(pdf)
        for page in pdf_reader.pages:
            text += page.extract_text()
    return text

def extract_text_from_docx(docx_docs):
    text = ""
    for docx in docx_docs:
        doc = Document(docx)
        for para in doc.paragraphs:
            text += para.text + "\n"
    return text

def extract_text_from_pptx(pptx_docs):
    text = ""
    for pptx in pptx_docs:
        presentation = Presentation(pptx)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    return text

def get_text_chunks(text):
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=10000, chunk_overlap=1000)
    chunks = text_splitter.split_text(text)
    return chunks

def get_vector_store(text_chunks):
    try:
        embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
        vector_store = FAISS.from_texts(text_chunks, embedding=embeddings)
        vector_store.save_local("faiss_index")
    except Exception as e:
        st.error(f"Error creating vector store: {e}")
        raise

def get_conversational_chain():
    prompt_template = """
    Answer the question only using the information found in the provided context. Focus on completeness and detail in the answer, ensuring no information is omitted if it's available in the context. Do not add any information, assumptions, or interpretations beyond what is directly stated in the context.

    If the question cannot be answered with the information in the context, respond only with: "Answer is not available in the context."

    Context:
    {context}

    Question:
    {question}

    Answer:
    """


    model = ChatGoogleGenerativeAI(model="models/gemini-pro", temperature=0.3)
    prompt = PromptTemplate(template=prompt_template, input_variables=["context", "question"])
    chain = load_qa_chain(model, chain_type="stuff", prompt=prompt)

    return chain

def get_summary(text):
    summary_prompt = """
    Summarize the following text into a concise summary:\n\n
    {context}
    """
    
    prompt = PromptTemplate(template=summary_prompt, input_variables=["context"])
    model = ChatGoogleGenerativeAI(model="models/gemini-pro", temperature=0.3)
    chain = LLMChain(prompt=prompt, llm=model)
    
    try:
        summary = chain.run({"context": text})
        return summary
    except Exception as e:
        st.error(f"Error generating summary: {e}")
        return "Error generating summary."

def user_input(user_question, language_code, context):
    try:
        embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
        new_db = FAISS.load_local("faiss_index", embeddings, allow_dangerous_deserialization=True)
        docs = new_db.similarity_search(user_question)
        chain = get_conversational_chain()
        context_text = " ".join([f"Q: {q} A: {a}" for q, a in context])
        response = chain({"input_documents": docs, "context": context_text, "question": user_question}, return_only_outputs=True)
        
        answer = response["output_text"]
        context.append((user_question, answer))
        
        st.subheader("Reply")
        st.write(answer)
        
        if language_code in LANGUAGES:
            translated_answer = translator.translate(answer, dest=language_code).text
            st.subheader("Translated Reply")
            st.markdown(f"{translated_answer}")
            
            tts = gTTS(text=translated_answer, lang=language_code)
            audio_stream = io.BytesIO()
            tts.write_to_fp(audio_stream)
            audio_stream.seek(0)
            
            st.audio(audio_stream, format='audio/mp3')
        else:
            st.error(f"Invalid language code: {language_code}")
    except Exception as e:
        st.error(f"Error during user input processing: {e}")

def main():
    st.set_page_config(page_title="DocuEase")
    st.title("DocuEase📄")

    st.header("Chat with Document")
    
    if "context" not in st.session_state:
        st.session_state.context = []
        st.session_state.chat_history = ""
    
    uploaded_files = st.file_uploader("Upload your PDF, DOCX, or PPTX files", accept_multiple_files=True)
    
    if uploaded_files:
        text = ""
        pdf_files = [file for file in uploaded_files if file.name.endswith('.pdf')]
        docx_files = [file for file in uploaded_files if file.name.endswith('.docx')]
        pptx_files = [file for file in uploaded_files if file.name.endswith('.pptx')]

        if pdf_files:
            text += extract_text_from_pdf(pdf_files)
        if docx_files:
            text += extract_text_from_docx(docx_files)
        if pptx_files:
            text += extract_text_from_pptx(pptx_files)

        if st.button("Submit & Process"):
            with st.spinner("Processing..."):
                try:
                    text_chunks = get_text_chunks(text)
                    get_vector_store(text_chunks)
                    st.success("Processing complete. You can now ask questions or summarize the document.")
                except Exception as e:
                    st.error(f"Error during processing: {e}")

        if st.button("Summarize Document"):
            with st.spinner("Generating summary..."):
                summary = get_summary(text)
                st.write("Summary: ", summary)

    user_question = st.text_input("Ask a Question from the Documents")
    language = st.selectbox("Select Language", options=list(LANGUAGES.keys()), format_func=lambda x: LANGUAGES[x])

    if user_question and language:
        user_input(user_question, language, st.session_state.context)
        st.session_state.chat_history += f"User: {user_question}\nBot: {st.session_state.context[-1][1]}\n\n"
        
    st.text_area("Chat History", value=st.session_state.chat_history, height=400)

if __name__ == "__main__":
    main()
